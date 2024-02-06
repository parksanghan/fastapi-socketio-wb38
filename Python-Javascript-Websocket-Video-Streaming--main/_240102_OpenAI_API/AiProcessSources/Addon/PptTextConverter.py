# -*- coding: cp949 -*-

from pptx import Presentation
import glob


def PptToText (root : str) :
    retList : list[str] = []    

    for eachfile in glob.glob(root):
        prs : Presentation = Presentation(eachfile)
    
        for slide in prs.slides:
            tStr : str = ""
            
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    tStr += shape.text;
            
            retList.append(tStr);
                

    return retList



